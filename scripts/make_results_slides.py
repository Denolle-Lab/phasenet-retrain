#!/usr/bin/env python3
"""
make_results_slides.py

Export all model evaluation results as a PowerPoint deck.
Slides:
  1. Title
  2. Overall (cross-domain, all distances)
  3. Local (<150 km)
  4. Regional (150-1500 km)
  5. Teleseismic (>1500 km)
  6. Notes on methodology

Run from repo root:
    conda activate surface
    python scripts/make_results_slides.py
"""

import numpy as np
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).parent.parent.resolve()
OUT_PPTX  = REPO_ROOT / "results" / "model_evaluation_results.pptx"

# ── Load & merge CSVs ──────────────────────────────────────────────────────────

eqt = pd.read_csv(REPO_ROOT / "results" / "eval_eqtransformer.csv")
ens = pd.read_csv(REPO_ROOT / "results" / "eval_ensemble_eqt.csv")
ens_row = ens[ens["weight"] == "eqt_ensemble_volpick_nc"]
combined = pd.concat(
    [eqt[eqt["weight"] != "eqt_ensemble_volpick_nc"], ens_row],
    ignore_index=True,
)
cross = combined[combined["split"] == "cross_domain"].copy()

# ── Model display names and grouping ──────────────────────────────────────────

MODEL_GROUPS = {
    "Ensemble": [
        "eqt_ensemble_volpick_nc",
    ],
    "EQTransformer": [
        "eqt_volpick",
        "eqt_original_nonconservative",
        "eqt_scedc",
        "eqt_instance",
        "eqt_ethz",
        "eqt_pnw",
        "eqt_original",
        "eqt_stead",
        "eqt_geofon",
        "eqt_iquique",
        "eqt_neic",
        "eqt_lendb",
    ],
    "PhaseNet": [
        "jma_wc_ft_global_v7",
        "jma_wc",
        "instance",
        "stead",
        "neic",
    ],
}

DISPLAY_NAMES = {
    "eqt_ensemble_volpick_nc":      "EQT Ensemble (volpick + orig_nc)",
    "eqt_volpick":                  "eqt_volpick",
    "eqt_original_nonconservative": "eqt_original_nc",
    "eqt_scedc":                    "eqt_scedc",
    "eqt_instance":                 "eqt_instance",
    "eqt_ethz":                     "eqt_ethz",
    "eqt_pnw":                      "eqt_pnw",
    "eqt_original":                 "eqt_original",
    "eqt_stead":                    "eqt_stead",
    "eqt_geofon":                   "eqt_geofon",
    "eqt_iquique":                  "eqt_iquique",
    "eqt_neic":                     "eqt_neic",
    "eqt_lendb":                    "eqt_lendb",
    "jma_wc_ft_global_v7":          "jma_wc_ft_v7 ★",
    "jma_wc":                       "jma_wc (parent)",
    "instance":                     "PhaseNet_instance",
    "stead":                        "PhaseNet_stead",
    "neic":                         "PhaseNet_neic",
}

METRICS = [
    ("p_mae_s",  "P-MAE (s)",  False),   # lower is better
    ("s_mae_s",  "S-MAE (s)",  False),
    ("p_recall", "P-Recall",   True),    # higher is better
    ("s_recall", "S-Recall",   True),
    ("mcc",      "MCC",        True),
]

DIST_BINS = [
    ("all",                     "All distances"),
    ("local (<150km)",          "Local  (<150 km)"),
    ("regional (150-1500km)",   "Regional  (150–1500 km)"),
    ("teleseismic (>1500km)",   "Teleseismic  (>1500 km)"),
]

# ── Colours ────────────────────────────────────────────────────────────────────

GROUP_HEADER_BG = {
    "Ensemble":     RGBColor(0xE7, 0x4C, 0x3C),   # red
    "EQTransformer":RGBColor(0x21, 0x6B, 0xD6),   # blue
    "PhaseNet":     RGBColor(0x27, 0xAE, 0x60),   # green
}
GROUP_ROW_BG = {
    "Ensemble":     RGBColor(0xFC, 0xE8, 0xE6),
    "EQTransformer":RGBColor(0xE8, 0xF2, 0xFF),
    "PhaseNet":     RGBColor(0xE9, 0xF7, 0xEF),
}
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
BEST_FG     = RGBColor(0x00, 0x00, 0x00)
BEST_BG     = RGBColor(0xFF, 0xD7, 0x00)   # gold highlight for best value


# ── Helper: pivot data for one dist_bin ───────────────────────────────────────

def get_table_data(dist_bin_key):
    """
    Returns ordered list of (group, weight, {metric: value_str}) dicts,
    plus a set of (col_idx, row_idx) that are 'best' in their metric column.
    """
    sub = cross[cross["dist_bin"] == dist_bin_key]
    rows = []
    for group, weights in MODEL_GROUPS.items():
        for w in weights:
            r = sub[sub["weight"] == w]
            vals = {}
            for col, _, _ in METRICS:
                if r.empty or col not in r.columns or np.isnan(r.iloc[0][col]):
                    vals[col] = "—"
                else:
                    v = r.iloc[0][col]
                    if col in ("p_mae_s", "s_mae_s", "p_outlier"):
                        vals[col] = f"{v:.3f}"
                    elif col in ("p_recall", "s_recall"):
                        vals[col] = f"{v:.3f}"
                    elif col == "mcc":
                        vals[col] = f"{v:.3f}"
                    else:
                        vals[col] = f"{v:.3f}"
            rows.append({"group": group, "weight": w, "vals": vals})

    # find best per metric column
    best = {}
    for col, _, higher in METRICS:
        numeric = []
        for i, row in enumerate(rows):
            try:
                numeric.append((float(row["vals"][col]), i))
            except (ValueError, KeyError):
                pass
        if numeric:
            best_val, best_i = (max if higher else min)(numeric, key=lambda x: x[0])
            best[col] = best_i
    return rows, best


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

BLANK = prs.slide_layouts[6]   # completely blank layout


def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=None, bg_color=None,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg_color
    return txb


def set_cell(cell, text, font_size=7.5, bold=False,
             bg_color=None, fg_color=None, align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = False
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if fg_color:
        run.font.color.rgb = fg_color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    # margins
    from pptx.util import Pt as Pt2
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    from pptx.oxml.ns import qn
    from lxml import etree
    for attr in ("marL", "marR", "marT", "marB"):
        tcPr.set(attr, str(Emu(Pt2(1.5))))


# ── Slide 1: Title ────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(BLANK)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

add_textbox(slide, "Seismic Phase Picker Benchmark",
            1, 2.5, 14, 1.2, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Cross-domain evaluation — P-MAE · S-MAE · P-Recall · S-Recall · MCC",
            1, 3.9, 14, 0.6, font_size=18, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, "Models: EQT Ensemble · 12× EQTransformer weights · 5× PhaseNet weights",
            1, 4.6, 14, 0.5, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(slide, "Benchmark: 6k-sample waveforms · threshold = 0.30 · search window ±5 s",
            1, 5.1, 14, 0.5, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# legend chips
for i, (group, col) in enumerate([
        ("Ensemble", GROUP_HEADER_BG["Ensemble"]),
        ("EQTransformer", GROUP_HEADER_BG["EQTransformer"]),
        ("PhaseNet", GROUP_HEADER_BG["PhaseNet"]),
]):
    x = 4.5 + i * 2.8
    tb = slide.shapes.add_textbox(Inches(x), Inches(6.2), Inches(2.4), Inches(0.4))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = group
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE
    tb.fill.solid(); tb.fill.fore_color.rgb = col


# ── Slides 2-5: one per distance bin ─────────────────────────────────────────

def make_table_slide(dist_key, dist_label):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

    # Title bar
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(15.4), Inches(0.55))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"Cross-domain Benchmark — {dist_label}"
    run.font.size = Pt(22); run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Subtitle
    add_textbox(slide, "threshold = 0.30 · ★ = best fine-tune · — = no picks in window · gold = best per column",
                0.3, 0.72, 15, 0.35, font_size=9,
                color=RGBColor(0x66, 0x66, 0x66), italic=True)

    rows_data, best = get_table_data(dist_key)

    N_ROWS = len(rows_data) + 2   # header + group-label rows (one per group)
    N_COLS = 1 + len(METRICS)     # model name + 5 metrics

    # Count group header rows
    groups_seen = []
    table_rows = []
    for rd in rows_data:
        if rd["group"] not in groups_seen:
            groups_seen.append(rd["group"])
            table_rows.append(("group_header", rd["group"]))
        table_rows.append(("data", rd))

    N_TABLE_ROWS = len(table_rows) + 1   # +1 for col header

    TBL_LEFT   = 0.25
    TBL_TOP    = 1.10
    TBL_WIDTH  = 15.5
    TBL_HEIGHT = 7.65

    tbl = slide.shapes.add_table(
        N_TABLE_ROWS, N_COLS,
        Inches(TBL_LEFT), Inches(TBL_TOP),
        Inches(TBL_WIDTH), Inches(TBL_HEIGHT),
    ).table

    # Column widths
    col_widths = [3.4] + [2.42] * len(METRICS)
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)

    # ── Header row ────────────────────────────────────────────────────────────
    HDR_BG = RGBColor(0x1A, 0x1A, 0x2E)
    set_cell(tbl.cell(0, 0), "Model", font_size=9, bold=True,
             bg_color=HDR_BG, fg_color=WHITE)
    for ci, (col, label, higher) in enumerate(METRICS, start=1):
        arrow = "▼ lower" if not higher else "▲ higher"
        set_cell(tbl.cell(0, ci), f"{label}\n({arrow})",
                 font_size=8.5, bold=True, bg_color=HDR_BG, fg_color=WHITE)

    # ── Data rows ─────────────────────────────────────────────────────────────
    data_row_idx = {}   # weight → table row index
    for ri, tr in enumerate(table_rows, start=1):
        kind = tr[0]
        if kind == "group_header":
            group = tr[1]
            bg = GROUP_HEADER_BG[group]
            merged_cell = tbl.cell(ri, 0)
            set_cell(merged_cell, f"── {group} ──",
                     font_size=8.5, bold=True, bg_color=bg, fg_color=WHITE,
                     align=PP_ALIGN.LEFT)
            for ci in range(1, N_COLS):
                set_cell(tbl.cell(ri, ci), "", bg_color=bg)
        else:
            rd = tr[1]
            group = rd["group"]
            weight = rd["weight"]
            data_row_idx[weight] = ri
            row_bg = GROUP_ROW_BG[group]
            name = DISPLAY_NAMES.get(weight, weight)
            set_cell(tbl.cell(ri, 0), name, font_size=8, bold=False,
                     bg_color=row_bg, fg_color=DARK_GRAY, align=PP_ALIGN.LEFT)
            for ci, (col, _, higher) in enumerate(METRICS, start=1):
                val_str = rd["vals"].get(col, "—")
                is_best = best.get(col) is not None and \
                          rows_data.index(rd) == best.get(col)
                bg = BEST_BG if is_best else row_bg
                fg = BEST_FG if is_best else DARK_GRAY
                set_cell(tbl.cell(ri, ci), val_str, font_size=8,
                         bold=is_best, bg_color=bg, fg_color=fg)

    return slide


for dist_key, dist_label in DIST_BINS:
    make_table_slide(dist_key, dist_label)


# ── Slide 6: Methodology notes ────────────────────────────────────────────────

slide = prs.slides.add_slide(BLANK)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)

add_textbox(slide, "Methodology Notes",
            0.4, 0.2, 15, 0.6, font_size=24, bold=True,
            color=RGBColor(0x1A, 0x1A, 0x2E))

notes = [
    ("Benchmark dataset",
     "~32k waveforms drawn from 11 seismic datasets (STEAD, INSTANCE, PNW, ETHZ, TXED, MLAAPDE, PISDL, CEED, VCSEIS, AQ2009GM, CWA). "
     "Stratified by distance bin (local / regional / teleseismic) and magnitude."),
    ("Waveform windows",
     "PhaseNet models evaluated on 3000-sample (30 s) windows. "
     "EQTransformer models evaluated on 6000-sample (60 s) windows extracted from raw data "
     "(noise-tiled where raw waveform too short), avoiding the zero-padding artifact."),
    ("Cross-domain split",
     "Traces from a model's training dataset are excluded from its cross-domain evaluation. "
     "EQT ensemble (volpick + original_nc) excludes STEAD traces since original_nc was trained on STEAD. "
     "jma_wc / jma_wc_ft_global_v7 show '—' here by design: their training manifest was assembled from "
     "every benchmark source dataset, so no benchmark trace's dataset is truly held out at this granularity. "
     "See results/event_leakage_audit.csv for the event-level (not dataset-level) independence check."),
    ("Metrics",
     "P/S-MAE: mean absolute timing error for detected picks. "
     "P/S-Recall: fraction of events where peak probability ≥ 0.30 within ±5 s of true arrival. "
     "MCC: Matthews Correlation Coefficient measuring P vs S phase discrimination."),
    ("Data leakage caveat",
     "EQT weights trained on STEAD / INSTANCE / ETHZ / PNW benefit from indirect distributional overlap "
     "with nearby datasets even after cross-domain exclusion. "
     "eqt_volpick (volcano-tectonic training data) has no overlap with benchmark datasets."),
    ("Ensemble",
     "Average of eqt_volpick (norm=peak) and eqt_original_nc (norm=std) probability curves. "
     "Each model's input batch normalised separately before inference."),
]

y = 0.95
for title, body in notes:
    add_textbox(slide, f"• {title}:", 0.5, y, 3.2, 0.35,
                font_size=10, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))
    add_textbox(slide, body, 3.7, y, 11.8, 0.45,
                font_size=9.5, color=DARK_GRAY)
    y += 0.58 if title != "Benchmark dataset" else 0.65


# ── Save ──────────────────────────────────────────────────────────────────────

prs.save(OUT_PPTX)
print(f"Saved → {OUT_PPTX}")
