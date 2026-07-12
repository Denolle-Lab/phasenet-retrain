"""
Generate phasenet_retraining_summary.pptx

Run from repo root:
    python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
import os
from lxml import etree

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x1a, 0x3a, 0x5c)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_LIGHT = RGBColor(0xC6, 0xEF, 0xCE)   # table row highlight
GREEN_TEXT  = RGBColor(0x27, 0x6F, 0x35)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY    = RGBColor(0xBF, 0xBF, 0xBF)
DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)
GOLD        = RGBColor(0xFF, 0xC0, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_NAME = "Calibri"

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]          # completely blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=Pt(18), bold=False, italic=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header(slide, title_text, subtitle_text=None):
    """Dark-blue header bar with white title (and optional subtitle)."""
    header_h = Inches(1.15)
    rect = add_rect(slide, 0, 0, SLIDE_W, header_h, fill_color=DARK_BLUE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.08),
                                          SLIDE_W - Inches(0.7), Inches(0.72))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.75),
                                            SLIDE_W - Inches(0.7), Inches(0.42))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.name = FONT_NAME
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xF0)

    return header_h


def add_bullet_body(slide, top, bullets, font_size=Pt(18), indent_size=Inches(0.35),
                    left=Inches(0.55), width=None, color=DARK_GRAY):
    if width is None:
        width = SLIDE_W - Inches(1.1)
    avail_h = SLIDE_H - top - Inches(0.2)
    txBox = slide.shapes.add_textbox(left, top, width, avail_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        level = item.get("level", 0)
        p.level = level
        p.space_before = Pt(4) if level == 0 else Pt(2)

        # bullet character
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        # indent
        buIndent = Inches(0.2 + level * 0.25)
        pPr.set('marL', str(int(buIndent)))
        pPr.set('indent', str(int(-Inches(0.2))))

        run = p.add_run()
        run.text = item["text"]
        run.font.name = FONT_NAME
        run.font.size = font_size if level == 0 else Pt(font_size.pt - 2)
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = item.get("color", color)


def add_image_centered(slide, img_path, top, available_height,
                        left=Inches(0.4), max_width=None):
    if max_width is None:
        max_width = SLIDE_W - Inches(0.8)
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    # fit within box
    if available_height * aspect <= max_width:
        h = available_height
        w = h * aspect
    else:
        w = max_width
        h = w / aspect
    l = left + (max_width - w) / 2
    slide.shapes.add_picture(img_path, l, top, w, h)
    return h


def add_caption(slide, text, top, font_size=Pt(13)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top,
                                      SLIDE_W - Inches(1.0), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)


# ---------------------------------------------------------------------------
# Table helper
# ---------------------------------------------------------------------------

def set_cell_bg(cell, rgb: RGBColor):
    """Set table cell background colour."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))


def set_cell_text(cell, text, font_size=Pt(14), bold=False, color=DARK_GRAY,
                  align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for run in p.runs:
        run.text = ""
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

def slide1_title(prs):
    slide = blank_slide(prs)
    # Full dark-blue background top third
    add_rect(slide, 0, 0, SLIDE_W, Inches(3.8), fill_color=DARK_BLUE)
    # Decorative bottom stripe
    add_rect(slide, 0, Inches(3.8), SLIDE_W, Inches(0.12), fill_color=RGBColor(0x2E, 0x75, 0xB6))

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), SLIDE_W - Inches(1.2), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PhaseNet Global Fine-Tuning"
    run.font.name = FONT_NAME; run.font.size = Pt(40); run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Progress Report"
    run2.font.name = FONT_NAME; run2.font.size = Pt(36); run2.font.bold = True
    run2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xF0)

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), SLIDE_W - Inches(1.2), Inches(0.8))
    tf2 = sub.text_frame; tf2.word_wrap = True
    p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Iterative retraining of jma_wc for cross-domain generalization"
    r3.font.name = FONT_NAME; r3.font.size = Pt(20); r3.font.color.rgb = RGBColor(0xB8, 0xD4, 0xF0)

    # Date & author
    info = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), SLIDE_W - Inches(1.2), Inches(1.2))
    tf3 = info.text_frame; tf3.word_wrap = True
    for txt, sz in [("June 12, 2026", Pt(20)), ("Akash Kharita · University of Washington", Pt(18))]:
        pp = tf3.add_paragraph() if tf3.paragraphs[0].runs else tf3.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run()
        rr.text = txt; rr.font.name = FONT_NAME; rr.font.size = sz
        rr.font.color.rgb = DARK_GRAY
        if txt != "June 12, 2026":
            info.text_frame.add_paragraph()
            p_new = info.text_frame.paragraphs[-1]
            p_new.alignment = PP_ALIGN.CENTER
            r_new = p_new.add_run()
            r_new.text = "Akash Kharita · University of Washington"
            r_new.font.name = FONT_NAME; r_new.font.size = Pt(18)
            r_new.font.color.rgb = DARK_GRAY
            break

    # rewrite cleanly
    info2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), SLIDE_W - Inches(1.2), Inches(1.5))
    tf4 = info2.text_frame
    for line, sz, clr in [
        ("June 12, 2026", Pt(20), DARK_GRAY),
        ("Akash Kharita  ·  University of Washington", Pt(17), RGBColor(0x70,0x70,0x70))
    ]:
        pp = tf4.paragraphs[0] if not tf4.paragraphs[0].runs else tf4.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run()
        rr.text = line; rr.font.name = FONT_NAME; rr.font.size = sz; rr.font.color.rgb = clr
    # remove the extra shapes
    # (info was already added — remove it from xml)
    sp = info._element
    sp.getparent().remove(sp)


def slide2_goal(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Project Goal")
    bullets = [
        {"text": "Improve PhaseNet's cross-domain P/S picking accuracy by fine-tuning on a diverse global dataset", "level": 0},
        {"text": "Base model: jma_wc (pretrained PhaseNet, seisbench)", "level": 1},
        {"text": ""},
        {"text": "Benchmark dataset: 31,992 held-out traces from global seismic networks", "level": 0},
        {"text": "Spans local, regional and teleseismic distance regimes", "level": 1},
        {"text": ""},
        {"text": "Key evaluation metrics:", "level": 0},
        {"text": "P-MAE (s) — mean absolute pick-time error", "level": 1},
        {"text": "P-Recall @ threshold 0.3 — detection completeness", "level": 1},
        {"text": "MCC — Matthews Correlation Coefficient (phase discrimination)", "level": 1},
        {"text": "P-Outlier fraction — picks with |error| > 1 s", "level": 1},
        {"text": ""},
        {"text": "Motivation: Out-of-the-box PhaseNet generalises poorly to stations / distance ranges outside its training domain", "level": 0, "bold": True},
    ]
    add_bullet_body(slide, top + Inches(0.25), bullets, font_size=Pt(18))


def slide3_dataset(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Training Dataset")

    # Stats row
    stats = [
        ("527,477", "Training traces"),
        ("125,218", "Validation traces"),
        ("275,660", "Test traces"),
        ("~67 GB", "RAM pre-loaded"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(0.95)
    box_top = top + Inches(0.18)
    gap = Inches(0.22)
    total_w = len(stats) * box_w + (len(stats)-1) * gap
    start_l = (SLIDE_W - total_w) / 2

    for i, (val, lbl) in enumerate(stats):
        l = start_l + i * (box_w + gap)
        add_rect(slide, l, box_top, box_w, box_h,
                 fill_color=DARK_BLUE,
                 line_color=None)
        # value
        vb = slide.shapes.add_textbox(l, box_top, box_w, Inches(0.55))
        tf = vb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.name = FONT_NAME; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
        # label
        lb = slide.shapes.add_textbox(l, box_top + Inches(0.50), box_w, Inches(0.45))
        tf2 = lb.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = lbl
        r2.font.name = FONT_NAME; r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xF0)

    # Network list
    net_top = box_top + box_h + Inches(0.12)
    net_box = slide.shapes.add_textbox(Inches(0.5), net_top, SLIDE_W - Inches(1.0), Inches(0.35))
    tf3 = net_box.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Networks: JMA · SCEDC · ETHZ · Iquique · GEOFON · PNW · NCEDC · NEIC · and more"
    r3.font.name = FONT_NAME; r3.font.size = Pt(15); r3.font.color.rgb = RGBColor(0x40, 0x70, 0xA0)

    img_path = "notebooks/training_spatial_distribution.png"
    img_top = net_top + Inches(0.38)
    avail_h = SLIDE_H - img_top - Inches(0.15)
    add_image_centered(slide, img_path, img_top, avail_h)


def slide4_approach(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Approach: Knowledge Distillation + Timing Loss")

    # Two column layout
    col_w = Inches(5.9)
    col_gap = Inches(0.5)
    col_top = top + Inches(0.25)
    col_h  = SLIDE_H - col_top - Inches(0.2)

    # Left column box
    add_rect(slide, Inches(0.3), col_top, col_w, col_h,
             fill_color=RGBColor(0xF0, 0xF5, 0xFB), line_color=DARK_BLUE, line_width=Pt(1))

    # Right column box
    add_rect(slide, Inches(0.3) + col_w + col_gap, col_top, col_w, col_h,
             fill_color=RGBColor(0xF5, 0xFA, 0xF5), line_color=RGBColor(0x27,0x6F,0x35), line_width=Pt(1))

    kd_bullets = [
        {"text": "Knowledge Distillation (KD)", "bold": True, "color": DARK_BLUE},
        {"text": "Student model learns from frozen jma_wc teacher"},
        {"text": "Prevents catastrophic forgetting of pretrained weights"},
        {"text": "KD loss: α · KL(student ∥ teacher) + cross-entropy"},
        {"text": "α controls how strongly the teacher constrains the student"},
        {"text": ""},
        {"text": "Hardware & Training Stack", "bold": True, "color": DARK_BLUE},
        {"text": "Optimizer: AdamW"},
        {"text": "Mixed precision: torch.amp (AMP)"},
        {"text": "Compilation: torch.compile for speed"},
        {"text": "Batch size: 1024  ·  GPU: RTX 3090 (24 GB)"},
    ]

    timing_bullets = [
        {"text": "Timing Loss (new in v4+)", "bold": True, "color": GREEN_TEXT},
        {"text": "Soft-argmax extracts a differentiable pick position"},
        {"text": "L1 loss between predicted and true arrival time (in samples)"},
        {"text": "Loss term: β · L1(soft_argmax(pred), true_pos)"},
        {"text": "β controls timing-loss weight vs cross-entropy + KD"},
        {"text": ""},
        {"text": "Hypothesis", "bold": True, "color": GREEN_TEXT},
        {"text": "KD alone prevents forgetting but doesn't explicitly sharpen timing"},
        {"text": "Adding a small β (0.01–0.1) should reduce P-MAE without hurting recall"},
        {"text": "v6 tests this at β=0.01 with v3's safe LR=5e-6"},
    ]

    def fill_col(left, bullets, color_key):
        txBox = slide.shapes.add_textbox(left + Inches(0.2), col_top + Inches(0.15),
                                          col_w - Inches(0.4), col_h - Inches(0.3))
        tf = txBox.text_frame; tf.word_wrap = True
        first = True
        for item in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = item["text"]
            r.font.name = FONT_NAME
            r.font.size = Pt(15) if item.get("bold") else Pt(14)
            r.font.bold = item.get("bold", False)
            r.font.color.rgb = item.get("color", DARK_GRAY)

    fill_col(Inches(0.3), kd_bullets, "blue")
    fill_col(Inches(0.3) + col_w + col_gap, timing_bullets, "green")

    # Column headers
    for l, txt, clr in [
        (Inches(0.3), "Knowledge Distillation", DARK_BLUE),
        (Inches(0.3) + col_w + col_gap, "Timing Loss", GREEN_TEXT),
    ]:
        hb = slide.shapes.add_textbox(l, col_top - Inches(0.0), col_w, Inches(0.0))
        # already covered by fill_col


def slide5_runs_table(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Training Runs Overview")

    rows_data = [
        ["Version", "LR", "KD α", "Timing β", "Scheduler", "Early Stop", "Epochs", "Best val_loss"],
        ["v3", "5×10⁻⁶", "0.3", "0", "ReduceLROnPlateau", "val_loss", "~96", "0.0655"],
        ["v4 ⚠", "1×10⁻⁴", "0.1", "0.1", "Cosine + warmup", "val_p_mae_s", "23", "0.0703 (early stop bug!)"],
        ["v5", "1×10⁻⁴", "0.1", "0.1", "Cosine + warmup", "val_loss", "123", "0.0655"],
        ["v6 ⟳", "5×10⁻⁶", "0.3", "0.01", "ReduceLROnPlateau", "val_loss", "in progress", "—"],
    ]

    n_rows = len(rows_data)
    n_cols = len(rows_data[0])

    tbl_top  = top + Inches(0.3)
    tbl_left = Inches(0.3)
    tbl_w    = SLIDE_W - Inches(0.6)
    tbl_h    = SLIDE_H - tbl_top - Inches(0.3)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    # Column widths (proportional)
    col_props = [0.07, 0.09, 0.07, 0.09, 0.17, 0.14, 0.10, 0.16]
    total_emu = int(tbl_w)
    for ci, prop in enumerate(col_props):
        table.columns[ci].width = int(total_emu * prop)

    # Row heights
    row_h = int(tbl_h / n_rows)
    for ri in range(n_rows):
        table.rows[ri].height = row_h

    row_colors = [DARK_BLUE, LIGHT_GRAY, RGBColor(0xFF,0xEB,0xEB), LIGHT_GRAY, RGBColor(0xEB,0xF5,0xEB)]
    row_text_colors = [WHITE, DARK_GRAY, DARK_GRAY, DARK_GRAY, DARK_GRAY]

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            bg = row_colors[ri]
            set_cell_bg(cell, bg)
            is_header = ri == 0
            bold = is_header or (ci == 0)
            fsz = Pt(14) if is_header else Pt(13)
            clr = row_text_colors[ri]
            if ri == 1:  # v3 best — subtle green tint
                set_cell_bg(cell, RGBColor(0xD9,0xEA,0xD3))
            set_cell_text(cell, cell_text, font_size=fsz, bold=bold, color=clr)

    # Note below table
    note = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - Inches(0.38),
                                     SLIDE_W - Inches(0.8), Inches(0.35))
    tf = note.text_frame; p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "★ v3 = best cross-domain performance so far    ⚠ v4 bug: ES monitored val_p_mae_s which peaked mid-warmup    ⟳ v6 still training"
    r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.italic = True
    r.font.color.rgb = RGBColor(0x50,0x50,0x50)


def slide6_v3_curves(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "v3 Training — Best Performing Model",
                     subtitle_text="LR = 5×10⁻⁶  |  KD α = 0.3  |  ReduceLROnPlateau  |  ~96 epochs")

    img_path = "results/finetune_jma_wc_global_v3/training_dashboard.png"
    img_top = top + Inches(0.12)
    avail_h = SLIDE_H - img_top - Inches(0.48)
    add_image_centered(slide, img_path, img_top, avail_h)

    cap_top = SLIDE_H - Inches(0.45)
    add_caption(slide, "Smooth convergence over ~96 epochs at LR = 5×10⁻⁶. Best val P-MAE = 0.37 s", cap_top)


def slide7_v4v5_curves(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "v4 & v5 — High-LR Cosine Schedule (Problematic Runs)",
                     subtitle_text="LR = 1×10⁻⁴  |  KD α = 0.1  |  Timing β = 0.1  |  Cosine + warmup")

    img_v4 = "results/finetune_jma_wc_global_v4/training_dashboard.png"
    img_v5 = "results/finetune_jma_wc_global_v5/training_dashboard.png"

    panel_w = (SLIDE_W - Inches(1.0)) / 2
    panel_l_v4 = Inches(0.3)
    panel_l_v5 = Inches(0.3) + panel_w + Inches(0.4)
    img_top = top + Inches(0.18)
    avail_h = SLIDE_H - img_top - Inches(0.7)

    from PIL import Image as PILImage

    def fit_image(path, left, avail_w, avail_h_):
        with PILImage.open(path) as im:
            iw, ih = im.size
        aspect = iw / ih
        if avail_h_ * aspect <= avail_w:
            h = avail_h_; w = h * aspect
        else:
            w = avail_w; h = w / aspect
        l = left + (avail_w - w) / 2
        slide.shapes.add_picture(path, l, img_top, w, h)

    fit_image(img_v4, panel_l_v4, panel_w, avail_h)
    fit_image(img_v5, panel_l_v5, panel_w, avail_h)

    # labels
    for l, lbl in [(panel_l_v4, "v4  (epochs: 23, stopped early due to ES bug)"),
                   (panel_l_v5, "v5  (epochs: 123, high-LR drift)")]:
        lb = slide.shapes.add_textbox(l, img_top - Inches(0.28), panel_w, Inches(0.28))
        tf = lb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lbl
        r.font.name = FONT_NAME; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = DARK_BLUE

    cap_top = SLIDE_H - Inches(0.52)
    add_caption(slide,
                "v4: stopped at epoch 3 (wrong ES monitor — val_p_mae_s peaked mid-warmup).  "
                "v5: ran 123 epochs but high LR caused drift from pretrained weights → recall collapses.",
                cap_top, font_size=Pt(13))


def slide8_v6_curves(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "v6 Training — Back to v3 LR Regime + Timing Loss",
                     subtitle_text="LR = 5×10⁻⁶  |  KD α = 0.3  |  Timing β = 0.01  |  ReduceLROnPlateau")

    img_path = "results/finetune_jma_wc_global_v6/training_dashboard.png"
    img_top = top + Inches(0.12)
    avail_h = SLIDE_H - img_top - Inches(0.48)
    add_image_centered(slide, img_path, img_top, avail_h)

    add_caption(slide,
                "v6: LR = 5×10⁻⁶, KD α = 0.3, timing β = 0.01.  ~60 epochs done — val loss 0.054 and still falling.",
                SLIDE_H - Inches(0.45))


def slide9_benchmark_dashboard(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Cross-Domain Benchmark — P-MAE Ranking (all distances)")

    img_path = "notebooks/step3_ft_dashboard.png"
    img_top = top + Inches(0.12)
    avail_h = SLIDE_H - img_top - Inches(0.5)
    add_image_centered(slide, img_path, img_top, avail_h)

    add_caption(slide,
                "v3 leads all fine-tuned variants. v4 / v5 regressed due to high learning rate.",
                SLIDE_H - Inches(0.45))


def slide10_recall(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "P-wave Detection Recall vs Threshold")

    img_path = "notebooks/step3_ft_recall_curves.png"
    img_top = top + Inches(0.18)
    avail_h = SLIDE_H - img_top - Inches(0.2)
    add_image_centered(slide, img_path, img_top, avail_h)


def slide11_distance(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Timing Error by Distance Bin")

    img_path = "notebooks/step3_ft_distance_bins.png"
    img_top = top + Inches(0.18)
    avail_h = SLIDE_H - img_top - Inches(0.2)
    add_image_centered(slide, img_path, img_top, avail_h)


def slide12_residuals(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Pick Residual Distributions")

    img_path = "notebooks/step3_ft_residuals.png"
    img_top = top + Inches(0.18)
    avail_h = SLIDE_H - img_top - Inches(0.2)
    add_image_centered(slide, img_path, img_top, avail_h)


def slide13_findings(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Key Findings")

    findings = [
        {"text": "v3 achieves best cross-domain performance: P-MAE = 0.368 s, P-Recall = 0.872",
         "bold": True, "color": GREEN_TEXT},
        {"text": "LR = 5×10⁻⁶ (v3) vs 1×10⁻⁴ (v4/v5): high LR is the primary driver of regression, not the timing loss"},
        {"text": "High LR drifts model far from pretrained jma_wc weights → P-Recall collapses to 0.30–0.41"},
        {"text": ""},
        {"text": "v4 bug confirmed: early stopping monitored val_p_mae_s, which peaked at epoch 3 (mid-warmup) → severely undertrained",
         "color": RGBColor(0xC0,0x00,0x00)},
        {"text": "v5 fix: switched to val_loss ES → model trained fully (123 ep) but still regressed vs v3 — confirming LR is the culprit"},
        {"text": ""},
        {"text": "Root cause: 1×10⁻⁴ destabilises KD — teacher signal too weak relative to gradient magnitude",
         "bold": True},
        {"text": ""},
        {"text": "v6 hypothesis: maintain v3 LR/KD regime (5×10⁻⁶, α=0.3), add small timing β=0.01",
         "color": GREEN_TEXT},
        {"text": "  → Test whether timing loss sharpens picks without sacrificing cross-domain recall", "level": 1},
        {"text": ""},
        {"text": "Baseline pretrained jma_wc still competitive (P-MAE 0.374 s, Recall 0.881) — any fine-tuned model must beat this"},
    ]
    add_bullet_body(slide, top + Inches(0.2), findings, font_size=Pt(17))


def slide14_summary_table(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Benchmark Summary — Cross-Domain, All Distances")

    rows_data = [
        ["Model", "P-MAE (s)", "P-Recall", "MCC", "P-Outlier"],
        ["jma_wc  (pretrained baseline)", "0.374", "0.881", "0.790", "0.071"],
        ["jma_wc_ft_global_v3  ★ BEST", "0.368", "0.872", "0.747", "0.071"],
        ["instance", "0.453", "0.841", "0.851", "0.090"],
        ["neic", "0.517", "0.554", "0.372", "0.104"],
        ["jma_wc_ft_global_v5", "0.921", "0.411", "0.196", "0.220"],
        ["jma_wc_ft_global_v4", "0.573", "0.296", "-0.064", "0.117"],
    ]

    # Best value per metric col (lower is better for MAE & Outlier, higher for Recall & MCC)
    best_per_col = {}
    for ci, direction in [(1, "low"), (2, "high"), (3, "high"), (4, "low")]:
        vals = []
        for ri in range(1, len(rows_data)):
            try:
                vals.append((float(rows_data[ri][ci]), ri))
            except ValueError:
                pass
        if vals:
            if direction == "low":
                best_per_col[ci] = min(vals, key=lambda x: x[0])[1]
            else:
                best_per_col[ci] = max(vals, key=lambda x: x[0])[1]

    n_rows = len(rows_data); n_cols = 5
    tbl_top  = top + Inches(0.28)
    tbl_left = Inches(0.3)
    tbl_w    = SLIDE_W - Inches(0.6)
    tbl_h    = SLIDE_H - tbl_top - Inches(0.55)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [0.38, 0.155, 0.155, 0.155, 0.155]
    for ci, p in enumerate(col_widths):
        table.columns[ci].width = int(int(tbl_w) * p)

    row_h = int(tbl_h / n_rows)
    for ri in range(n_rows):
        table.rows[ri].height = row_h

    row_bgs = {
        0: DARK_BLUE,
        1: LIGHT_GRAY,   # baseline
        2: GREEN_LIGHT,  # v3 best
        3: LIGHT_GRAY,
        4: LIGHT_GRAY,
        5: RGBColor(0xFF,0xEB,0xEB),  # bad
        6: RGBColor(0xFF,0xEB,0xEB),  # bad
    }

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            bg = row_bgs.get(ri, LIGHT_GRAY)
            set_cell_bg(cell, bg)
            is_header = ri == 0
            is_best_cell = (ci in best_per_col and best_per_col[ci] == ri)
            txt_color = WHITE if is_header else DARK_GRAY
            if ri == 2:
                txt_color = GREEN_TEXT
            if ri in (5, 6):
                txt_color = RGBColor(0x99, 0x00, 0x00)
            bold = is_header or (ci == 0) or is_best_cell
            fsz = Pt(15) if is_header else Pt(14)
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, cell_text, font_size=fsz, bold=bold, color=txt_color, align=align)

    # Legend
    leg = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - Inches(0.48),
                                    SLIDE_W - Inches(0.8), Inches(0.40))
    tf = leg.text_frame; p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("★ Best fine-tuned model  |  Bold values = best in column  |  "
              "Green row = v3 (best FT)  |  Red rows = high-LR runs (regressed)  |  "
              "Lower P-MAE & P-Outlier is better;  Higher Recall & MCC is better")
    r.font.name = FONT_NAME; r.font.size = Pt(11); r.font.italic = True
    r.font.color.rgb = RGBColor(0x50,0x50,0x50)


def slide15_next_steps(prs):
    slide = blank_slide(prs)
    top = add_header(slide, "Next Steps")

    bullets = [
        {"text": "Immediate", "bold": True, "color": DARK_BLUE},
        {"text": "Evaluate v6 on the benchmark when training finishes (~2 hours)", "level": 0},
        {"text": ""},
        {"text": "Decision tree based on v6 results:", "bold": True, "color": DARK_BLUE},
        {"text": "If v6 beats v3 → timing loss at β = 0.01 is beneficial → explore β = 0.05", "level": 0,
         "color": GREEN_TEXT},
        {"text": "If v6 matches v3 → timing loss is neutral → adopt v3 as production model", "level": 0,
         "color": RGBColor(0xC0,0x70,0x00)},
        {"text": "If v6 regresses → timing loss hurts even at β = 0.01 → drop it entirely", "level": 0,
         "color": RGBColor(0xC0,0x00,0x00)},
        {"text": ""},
        {"text": "Deployment", "bold": True, "color": DARK_BLUE},
        {"text": "Export best model checkpoint to SeisBench format for community sharing", "level": 0},
        {"text": "Benchmark against additional held-out datasets (Ridgecrest, STEAD)", "level": 0},
        {"text": ""},
        {"text": "Longer-term", "bold": True, "color": DARK_BLUE},
        {"text": "Consider curriculum training: start with local events, add teleseismic progressively", "level": 0},
        {"text": "Explore S-wave timing loss (currently only P is targeted)", "level": 0},
        {"text": "Investigate ensemble of v3 + v6 for improved robustness", "level": 0},
    ]
    add_bullet_body(slide, top + Inches(0.22), bullets, font_size=Pt(18))


# ---------------------------------------------------------------------------
# Assemble
# ---------------------------------------------------------------------------

def build():
    prs = new_prs()

    slide1_title(prs)
    slide2_goal(prs)
    slide3_dataset(prs)
    slide4_approach(prs)
    slide5_runs_table(prs)
    slide6_v3_curves(prs)
    slide7_v4v5_curves(prs)
    slide8_v6_curves(prs)
    slide9_benchmark_dashboard(prs)
    slide10_recall(prs)
    slide11_distance(prs)
    slide12_residuals(prs)
    slide13_findings(prs)
    slide14_summary_table(prs)
    slide15_next_steps(prs)

    out = "phasenet_retraining_summary.pptx"
    prs.save(out)
    print(f"Saved → {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
